"""Bundle all PNG figures from a reports/figures/<version> directory into a pptx.

Usage:
    python scripts/figures_to_pptx.py [version]

Defaults to v0.50.1. Each PNG becomes one slide with the filename as the title,
sized to fit within the slide while preserving aspect ratio.
"""

import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Inches, Pt

VERSION = sys.argv[1] if len(sys.argv) > 1 else "v0.50.1"

REPO_ROOT = Path(__file__).resolve().parents[1]
FIG_DIR = REPO_ROOT / "reports" / "figures" / VERSION
COMPARISONS_DIR = REPO_ROOT / "reports" / "figures" / "comparisons"
OUT_PATH = FIG_DIR / f"gelos_lc_figures_{VERSION}.pptx"


def main() -> None:
    if not FIG_DIR.is_dir():
        sys.exit(f"figures dir not found: {FIG_DIR}")

    pngs = sorted(FIG_DIR.glob("*.png"))
    pngs += sorted(COMPARISONS_DIR.rglob("*.png"))
    if not pngs:
        sys.exit(f"no PNGs in {FIG_DIR}")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide_w, slide_h = prs.slide_width, prs.slide_height
    title_h = Inches(0.5)
    margin = Inches(0.2)
    avail_w = slide_w - 2 * margin
    avail_h = slide_h - title_h - 2 * margin

    for png in pngs:
        slide = prs.slides.add_slide(blank)

        tb = slide.shapes.add_textbox(margin, margin, avail_w, title_h)
        p = tb.text_frame.paragraphs[0]
        p.text = png.stem
        p.font.size = Pt(14)
        p.font.bold = True

        with Image.open(png) as im:
            iw, ih = im.size
        aspect = iw / ih
        if avail_w / avail_h > aspect:
            h = avail_h
            w = Emu(int(int(avail_h) * aspect))
        else:
            w = avail_w
            h = Emu(int(int(avail_w) / aspect))
        left = Emu(int((slide_w - w) / 2))
        top = title_h + margin + Emu(int((avail_h - h) / 2))
        slide.shapes.add_picture(str(png), left, top, width=w, height=h)

    prs.save(OUT_PATH)
    print(f"wrote {len(pngs)} slides -> {OUT_PATH}")


if __name__ == "__main__":
    main()
